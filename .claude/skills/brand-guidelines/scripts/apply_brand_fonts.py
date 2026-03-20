#!/usr/bin/env python3
"""
Apply brand fonts to a PPTX presentation.

This script post-processes a PPTX file to apply brand typography using
font-name-based mapping (not size-based):

- Arial Black → Bebas Neue (display headings)
- Arial (bold+italic) → Roboto SemiBold (clear bold, keep italic)
- Arial (bold only) → Roboto Medium (clear bold)
- Arial (normal, ≤18pt) → Roboto Medium
- Arial (normal, >18pt) → Roboto Light
- Courier/Mono → JetBrains Mono

Usage:
    python apply_brand_fonts.py input.pptx [output.pptx]

Dependencies:
    pip install python-pptx
"""

import sys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE

FONT_DISPLAY = "Bebas Neue"
FONT_BODY_LIGHT = "Roboto Light"
FONT_BODY_MEDIUM = "Roboto Medium"
FONT_BODY_SEMIBOLD = "Roboto SemiBold"
FONT_CODE = "JetBrains Mono"

# Width expansion factor for text boxes
WIDTH_EXPANSION = 1.05


def apply_brand_fonts(input_path, output_path):
    """Apply brand fonts to all text in the presentation."""
    prs = Presentation(input_path)

    stats = {"display": 0, "light": 0, "medium": 0, "semibold": 0, "code": 0}

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue

            tf = shape.text_frame
            has_display = False

            for para in tf.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue

                    current_font = run.font.name or ""
                    font_size = run.font.size
                    size_pt = font_size.pt if font_size else 14
                    is_bold = run.font.bold
                    is_italic = run.font.italic

                    # 1. Monospace/code fonts
                    if "Courier" in current_font or "Mono" in current_font:
                        run.font.name = FONT_CODE
                        stats["code"] += 1
                        continue

                    # 2. Arial Black → Bebas Neue (display headings)
                    if "Arial Black" in current_font:
                        run.font.name = FONT_DISPLAY
                        stats["display"] += 1
                        has_display = True
                        continue

                    # 3. Everything else (Arial, sans-serif, etc.) → Roboto variants
                    if is_bold and is_italic:
                        run.font.name = FONT_BODY_SEMIBOLD
                        run.font.bold = False  # weight is in the font name
                        # keep italic
                        stats["semibold"] += 1
                    elif is_bold and not is_italic:
                        run.font.name = FONT_BODY_MEDIUM
                        run.font.bold = False  # weight is in the font name
                        stats["medium"] += 1
                    elif size_pt <= 18:
                        run.font.name = FONT_BODY_MEDIUM
                        stats["medium"] += 1
                    else:
                        run.font.name = FONT_BODY_LIGHT
                        stats["light"] += 1

            # Adjust text frame for display headings
            if has_display:
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Expand text box width slightly
            if shape.width:
                try:
                    new_width = int(shape.width * WIDTH_EXPANSION)
                    width_delta = new_width - shape.width
                    shape.left = shape.left - int(width_delta / 2)
                    shape.width = new_width
                except Exception:
                    pass

    prs.save(output_path)

    print(f"Brand fonts applied:")
    print(f"  Display ({FONT_DISPLAY}): {stats['display']}")
    print(f"  Body Light ({FONT_BODY_LIGHT}): {stats['light']}")
    print(f"  Body Medium ({FONT_BODY_MEDIUM}): {stats['medium']}")
    print(f"  Body SemiBold ({FONT_BODY_SEMIBOLD}): {stats['semibold']}")
    print(f"  Code ({FONT_CODE}): {stats['code']}")
    print(f"\nSaved to: {output_path}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python apply_brand_fonts.py input.pptx [output.pptx]")
        print("\nIf output not specified, overwrites input file.")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else input_path

    apply_brand_fonts(input_path, output_path)


if __name__ == "__main__":
    main()
